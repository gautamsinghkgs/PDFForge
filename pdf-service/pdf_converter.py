import os
import logging
from pathlib import Path
from pdf2docx import Converter
import pdfplumber
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from docx import Document
import tempfile

logger = logging.getLogger(__name__)

async def pdf_to_word(input_path: Path, output_path: Path):
    """Convert PDF to DOCX using pdf2docx - best quality conversion"""
    try:
        logger.info(f"Starting PDF to Word conversion: {input_path}")
        logger.info(f"Input file exists: {input_path.exists()}")
        
        if not input_path.exists():
            raise Exception(f"Input file does not exist: {input_path}")
        
        # Ensure output directory exists
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # Use pdf2docx for best quality conversion
        logger.info("Initializing pdf2docx converter...")
        cv = Converter(str(input_path))
        
        logger.info("Starting conversion process...")
        cv.convert(str(output_path), start=0, end=None)
        cv.close()
        
        logger.info("Conversion process completed")
        
        # Wait for file system to sync
        import time
        time.sleep(0.5)
        
        # Verify file was created
        if output_path.exists():
            size = output_path.stat().st_size
            logger.info(f"PDF to Word conversion completed: {output_path}")
            logger.info(f"Output file size: {size} bytes")
            return
        else:
            # Check if file was created with different name in same directory
            temp_dir = output_path.parent
            docx_files = list(temp_dir.glob("*.docx"))
            logger.info(f"Found DOCX files: {[f.name for f in docx_files]}")
            
            if docx_files:
                # Use the most recent file
                latest_docx = max(docx_files, key=lambda x: x.stat().st_mtime)
                logger.info(f"Using most recent DOCX: {latest_docx}")
                
                # Copy to expected location
                import shutil
                shutil.copy2(latest_docx, output_path)
                logger.info(f"Copied to expected location: {output_path}")
                
                if output_path.exists():
                    size = output_path.stat().st_size
                    logger.info(f"PDF to Word conversion completed (with copy): {output_path}")
                    logger.info(f"Output file size: {size} bytes")
                    return
            
            raise Exception(f"Output file not created at {output_path}")
        
    except Exception as e:
        logger.error(f"pdf2docx conversion failed: {str(e)}")
        import traceback
        logger.error(f"Full traceback: {traceback.format_exc()}")
        # Fallback to text-based conversion
        await pdf_to_word_fallback(input_path, output_path)

async def pdf_to_word_fallback(input_path: Path, output_path: Path):
    """Improved PDF to Word conversion preserving layout using pdfplumber"""
    try:
        logger.info(f"Starting layout-preserving PDF to Word conversion: {input_path}")
        
        # Create Word document with A4 page size
        doc = Document()
        
        # Set default font
        from docx.shared import Pt, Inches, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        # Configure document sections for A4
        sections = doc.sections
        for section in sections:
            section.page_height = Inches(11.69)
            section.page_width = Inches(8.27)
            section.left_margin = Inches(0.75)
            section.right_margin = Inches(0.75)
            section.top_margin = Inches(0.75)
            section.bottom_margin = Inches(0.75)
        
        with pdfplumber.open(input_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                logger.info(f"Processing page {page_num + 1}")
                
                # Extract words with positions
                words = page.extract_words(
                    keep_blank_chars=True,
                    x_tolerance=3,
                    y_tolerance=3
                )
                
                if not words:
                    # Fallback to simple text extraction if no words found
                    text = page.extract_text()
                    if text:
                        paragraphs = text.split('\n')
                        for para_text in paragraphs:
                            if para_text.strip():
                                p = doc.add_paragraph(para_text.strip())
                                p.paragraph_format.space_after = Pt(6)
                    continue
                
                # Group words by lines (y-coordinate)
                lines = {}
                for word in words:
                    y_pos = round(word['top'], 1)  # Group by rounded y position
                    if y_pos not in lines:
                        lines[y_pos] = []
                    lines[y_pos].append(word)
                
                # Sort lines by y position (top to bottom)
                sorted_y_positions = sorted(lines.keys())
                
                for y_pos in sorted_y_positions:
                    line_words = lines[y_pos]
                    # Sort words by x position (left to right)
                    line_words.sort(key=lambda w: w['x0'])
                    
                    # Combine words into line text
                    line_text = ' '.join(word['text'] for word in line_words)
                    
                    if line_text.strip():
                        # Create paragraph with original spacing
                        p = doc.add_paragraph(line_text)
                        
                        # Set paragraph formatting to preserve spacing
                        p.paragraph_format.space_after = Pt(3)
                        p.paragraph_format.space_before = Pt(0)
                        p.paragraph_format.line_spacing = 1.15
                        
                        # Check if this might be a heading (based on font size and position)
                        avg_font_size = sum(w.get('size', 12) for w in line_words) / len(line_words) if line_words else 12
                        
                        if avg_font_size > 14:
                            # Apply heading style
                            run = p.runs[0] if p.runs else p.add_run(line_text)
                            run.font.size = Pt(min(avg_font_size, 16))
                            run.font.bold = True
                        else:
                            # Normal text
                            run = p.runs[0] if p.runs else p.add_run(line_text)
                            run.font.size = Pt(11)
                
                # Add page break after each page except the last
                if page_num < len(pdf.pages) - 1:
                    doc.add_page_break()
        
        # Save the document
        doc.save(str(output_path))
        logger.info(f"Layout-preserving PDF to Word conversion completed: {output_path}")
        
    except Exception as e:
        logger.error(f"Layout-preserving conversion failed: {str(e)}")
        import traceback
        logger.error(traceback.format_exc())
        # Fallback to basic conversion
        await basic_pdf_to_word_fallback(input_path, output_path)

async def basic_pdf_to_word_fallback(input_path: Path, output_path: Path):
    """Basic fallback when layout preservation fails"""
    try:
        logger.info("Using basic fallback PDF to Word conversion")
        
        doc = Document()
        text_content = []
        
        with pdfplumber.open(input_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_content.append(text)
        
        for page_text in text_content:
            paragraphs = page_text.split('\n')
            for para in paragraphs:
                if para.strip():
                    doc.add_paragraph(para.strip())
            doc.add_page_break()
        
        doc.save(str(output_path))
        logger.info(f"Basic fallback conversion completed: {output_path}")
        
    except Exception as e:
        logger.error(f"Basic fallback conversion failed: {str(e)}")
        raise Exception(f"PDF to Word conversion failed: {str(e)}")

async def pdf_to_excel(input_path: Path, output_path: Path):
    """Convert PDF to XLSX using xlsxwriter - most reliable method"""
    try:
        logger.info(f"Starting reliable PDF to Excel conversion: {input_path}")
        
        import xlsxwriter
        
        # Create workbook with xlsxwriter (most reliable)
        workbook = xlsxwriter.Workbook(str(output_path))
        worksheet = workbook.add_worksheet('PDF Content')
        
        # Define formats
        header_format = workbook.add_format({
            'bold': True,
            'font_size': 14,
            'bg_color': '#4472C4',
            'font_color': 'white',
            'border': 1
        })
        
        text_format = workbook.add_format({
            'font_size': 11,
            'text_wrap': True,
            'valign': 'top'
        })
        
        row = 0
        
        with pdfplumber.open(input_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                # Write page header
                worksheet.write(row, 0, f"Page {page_num + 1}", header_format)
                row += 1
                
                # Extract text
                text = page.extract_text()
                if text:
                    lines = text.split('\n')
                    for line in lines:
                        if line.strip():
                            worksheet.write(row, 0, line.strip(), text_format)
                            row += 1
                
                # Add blank row between pages
                row += 1
        
        # Set column width
        worksheet.set_column('A:A', 80)
        
        # Close workbook (this saves the file)
        workbook.close()
        
        logger.info(f"Reliable PDF to Excel conversion completed: {output_path}")
        
    except Exception as e:
        logger.error(f"PDF to Excel conversion failed: {str(e)}")
        import traceback
        logger.error(traceback.format_exc())
        raise Exception(f"PDF to Excel conversion failed: {str(e)}")

async def pdf_to_csv(input_path: Path, output_path: Path):
    """Convert PDF to CSV - 100% reliable method"""
    try:
        logger.info(f"Starting PDF to CSV conversion: {input_path}")
        
        import csv
        
        with open(output_path, 'w', newline='', encoding='utf-8') as csvfile:
            writer = csv.writer(csvfile)
            writer.writerow(['Page', 'Content'])  # Header
            
            with pdfplumber.open(input_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    if text:
                        for line in text.split('\n'):
                            if line.strip():
                                writer.writerow([f"Page {page_num + 1}", line.strip()])
                    else:
                        writer.writerow([f"Page {page_num + 1}", "(No text content)"])
        
        logger.info(f"PDF to CSV conversion completed: {output_path}")
        
    except Exception as e:
        logger.error(f"PDF to CSV conversion failed: {str(e)}")
        import traceback
        logger.error(traceback.format_exc())
        raise Exception(f"PDF to CSV conversion failed: {str(e)}")

async def pdf_to_powerpoint(input_path: Path, output_path: Path):
    """Convert PDF to PPTX with A4 size and structure preservation"""
    try:
        logger.info(f"Starting PDF to PowerPoint conversion with A4 layout: {input_path}")
        
        # Create presentation with custom A4 size
        prs = Presentation()
        
        # A4 dimensions: 8.27 x 11.69 inches
        A4_WIDTH = Inches(8.27)
        A4_HEIGHT = Inches(11.69)
        
        # Set slide dimensions to A4
        prs.slide_width = A4_WIDTH
        prs.slide_height = A4_HEIGHT
        
        # Create custom blank layout
        blank_layout = prs.slide_layouts[6]  # Blank layout
        
        with pdfplumber.open(input_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                # Add slide
                slide = prs.slides.add_slide(blank_layout)
                
                # Extract words with positions for layout preservation
                words = page.extract_words(
                    keep_blank_chars=True,
                    x_tolerance=3,
                    y_tolerance=3
                )
                
                if words:
                    # Group words by lines (y-coordinate)
                    lines = {}
                    for word in words:
                        y_pos = round(word['top'], 1)
                        if y_pos not in lines:
                            lines[y_pos] = []
                        lines[y_pos].append(word)
                    
                    # Track previous text box position to avoid overlapping
                    prev_bottom = 0.5  # Start margin
                    min_spacing = Inches(0.15)  # Minimum spacing between lines
                    
                    # Sort by vertical position
                    for y_pos in sorted(lines.keys()):
                        line_words = sorted(lines[y_pos], key=lambda w: w['x0'])
                        line_text = ' '.join(w['text'] for w in line_words)
                        
                        if line_text.strip():
                            # Calculate position relative to A4 slide
                            pdf_height = page.height
                            pdf_width = page.width
                            
                            # Calculate scaled position
                            top_inches = 0.5 + (y_pos / pdf_height) * 10.69 if pdf_height > 0 else 0.5
                            
                            # Ensure minimum spacing from previous text box
                            if top_inches < prev_bottom + min_spacing.inches:
                                top_inches = prev_bottom + min_spacing.inches
                            
                            # Check if we're going beyond page
                            if top_inches > 11.0:
                                break  # Stop adding text if we're at bottom of page
                            
                            left = Inches(0.5 + (line_words[0]['x0'] / pdf_width) * 7.27) if pdf_width > 0 else Inches(0.5)
                            top = Inches(top_inches)
                            
                            # Calculate text box height based on content length
                            text_length = len(line_text)
                            estimated_lines = max(1, text_length // 60)  # Estimate lines (60 chars per line)
                            box_height = Inches(0.25 * estimated_lines + 0.1)  # Minimum 0.35 inches
                            
                            # Check average font size for this line
                            avg_font_size = sum(w.get('size', 12) for w in line_words) / len(line_words) if line_words else 12
                            scaled_font_size = min(max(avg_font_size, 10), 16)  # Keep reasonable size
                            
                            # Add text box
                            textbox = slide.shapes.add_textbox(
                                left, top, Inches(7), box_height
                            )
                            text_frame = textbox.text_frame
                            text_frame.text = line_text
                            
                            # Format text
                            paragraph = text_frame.paragraphs[0]
                            paragraph.font.size = Pt(scaled_font_size)
                            paragraph.font.name = 'Arial'
                            
                            # Make headings bold
                            if avg_font_size > 14 or line_text.isupper():
                                paragraph.font.bold = True
                            
                            text_frame.word_wrap = True
                            
                            # Update previous bottom position
                            prev_bottom = top_inches + box_height.inches
                
                else:
                    # Fallback to basic text extraction
                    text = page.extract_text()
                    if text:
                        # Add centered text box
                        textbox = slide.shapes.add_textbox(
                            Inches(0.5), Inches(0.5), Inches(7.27), Inches(10.69)
                        )
                        text_frame = textbox.text_frame
                        text_frame.text = text
                        paragraph = text_frame.paragraphs[0]
                        paragraph.font.size = Pt(12)
                        paragraph.font.name = 'Arial'
                        text_frame.word_wrap = True
                
                # Add page number at bottom
                page_num_box = slide.shapes.add_textbox(
                    Inches(3.5), Inches(11), Inches(1.27), Inches(0.5)
                )
                page_num_frame = page_num_box.text_frame
                page_num_frame.text = f"Page {page_num + 1}"
                page_num_para = page_num_frame.paragraphs[0]
                page_num_para.font.size = Pt(10)
                page_num_para.font.italic = True
                page_num_para.alignment = PP_ALIGN.CENTER
        
        # Save presentation
        prs.save(str(output_path))
        logger.info(f"PDF to PowerPoint conversion completed with A4 layout: {output_path}")
        
    except Exception as e:
        logger.error(f"PDF to PowerPoint conversion failed: {str(e)}")
        import traceback
        logger.error(traceback.format_exc())
        raise Exception(f"PDF to PowerPoint conversion failed: {str(e)}")

async def word_to_pdf(input_path: Path, output_path: Path):
    """Convert DOCX to PDF - placeholder for future implementation"""
    raise NotImplementedError("Word to PDF conversion not implemented yet")

async def excel_to_pdf(input_path: Path, output_path: Path):
    """Convert XLSX to PDF - placeholder for future implementation"""
    raise NotImplementedError("Excel to PDF conversion not implemented yet")

async def powerpoint_to_pdf(input_path: Path, output_path: Path):
    """Convert PPTX to PDF - placeholder for future implementation"""
    raise NotImplementedError("PowerPoint to PDF conversion not implemented yet")
